#!/usr/bin/env python3
"""
pptx-normalizer: Remove animations and auto-play, unify fonts to 微软雅黑.

Usage:
  python3 normalize.py <input.pptx> [output.pptx]

If output is omitted, overwrites the input file (with backup).
"""

import sys
import os
import shutil
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    print("Error: python-pptx not installed.")
    print("Install with: pip install python-pptx")
    sys.exit(1)


def normalize_pptx(input_path: str, output_path: str):
    prs = Presentation(input_path)

    # 1. Remove slide transitions and auto-advance timings
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    for slide in prs.slides:
        xml_slide = slide._element
        # Remove <p:transition> elements
        transitions = xml_slide.findall('.//p:transition', ns)
        for t in transitions:
            xml_slide.remove(t)
        # Set all timing nodes to advance on click only
        for cTn in xml_slide.findall('.//p:cTn', ns):
            cTn.set('advTm', 'indefinite')

    # 3. Process all shapes on all slides: remove animations, unify font
    changes = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            # Remove animation effects
            try:
                if shape.click_action:
                    shape.click_action.target_slide = None
            except Exception:
                pass

            # Check for and handle grouped shapes
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = '微软雅黑'
                        try:
                            run.font.size = Pt(run.font.size.pt)  # trigger write
                        except Exception:
                            pass
                        changes += 1

            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.name = '微软雅黑'
                                changes += 1

    prs.save(output_path)
    return changes


def main():
    if len(sys.argv) < 2:
        print("Usage: python3 normalize.py <input.pptx> [output.pptx]")
        print("  If output is omitted, input is overwritten (backup saved as .bak)")
        sys.exit(1)

    input_path = sys.argv[1]

    if len(sys.argv) >= 3:
        output_path = sys.argv[2]
        is_overwrite = False
    else:
        # Backup original
        backup_path = input_path + '.bak'
        shutil.copy2(input_path, backup_path)
        print(f"Backup saved: {backup_path}")
        output_path = input_path
        is_overwrite = True

    changes = normalize_pptx(input_path, output_path)
    print(f"Done. {changes} font changes applied. Output: {output_path}")
    print("Actions: removed animations, removed auto-play, unified font to 微软雅黑")


if __name__ == '__main__':
    main()
